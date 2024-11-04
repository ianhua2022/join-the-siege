from werkzeug.datastructures import FileStorage
from src.validator import categorize_file
from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import io
from PIL import Image
import easyocr
import numpy as np
from src.classifiers.civil import civil_classifier
from src.classifiers.finance import finance_classifier

# Supported industries, can be further expanded
SUPPORTED_INDUSTRIES = ['civil', 'finance']

def get_file_content(file: FileStorage) -> str:
    """
    Extracts text content from a file.

    Args:
        file (FileStorage): The file to extract content from, provided as a FileStorage object

    Returns:
        str: The extracted text content of the file
    """
    # Reset file pointer to beginning
    file.stream.seek(0)
    
    file_type = categorize_file(file)
    try:
        if file_type == "document":
                if file.content_type == 'application/pdf':
                    reader = PdfReader(file.stream)
                    content = ""
                    for page in reader.pages:
                        content += page.extract_text()
                    return content
                    
                elif file.content_type in ['application/msword', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document']:
                    # Create a copy of the stream for docx
                    file_copy = io.BytesIO(file.stream.read())
                    document = Document(file_copy)
                    return "\n".join(paragraph.text for paragraph in document.paragraphs)
                    
                elif file.content_type in ['application/vnd.ms-excel', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet']:
                    # Create a copy of the stream for Excel
                    file_copy = io.BytesIO(file.stream.read())
                    workbook = load_workbook(filename=file_copy)
                    content = ""
                    for sheet in workbook.worksheets:
                        for row in sheet.rows:
                            content += " ".join(str(cell.value or "") for cell in row) + "\n"
                    return content
                    
                elif file.content_type in ['application/vnd.ms-powerpoint', 'application/vnd.openxmlformats-officedocument.presentationml.presentation']:
                    # Create a copy of the stream for PowerPoint
                    file_copy = io.BytesIO(file.stream.read())
                    presentation = Presentation(file_copy)
                    content = ""
                    for slide in presentation.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                content += shape.text + "\n"
                    return content
                    
                elif file.content_type == 'text/plain':
                    return file.stream.read().decode('utf-8')
                    
                else:
                    return ""
                
        elif file_type == "image":
            # Use EasyOCR instead of pytesseract
            reader = easyocr.Reader(['en'])  # Initialize once for English
            file_copy = io.BytesIO(file.stream.read())
            image = Image.open(file_copy)
            result = reader.readtext(np.array(image))
            return ' '.join(item[1] for item in result)
            
        else:
            return ""
    except Exception as e:
        print(f"Error processing file: {str(e)}")
        return ""
        

def classify_industry(file_content: str) -> str:
    """
    Determines the industry classification of a document based on its content.

    Args:
        file_content (str): The text content of the file to classify

    Returns:
        str: The predicted industry classification. Can be one of:
            - 'civil': If civil industry keywords are found
            - 'finance': If finance industry keywords are found 
            - 'unknown': If no industry-specific keywords are found

    The function works by:
    1. Converting the content to lowercase and splitting into words
    2. Checking each word against industry-specific keyword lists
    3. Returning the first matching industry, or 'unknown' if no matches
    
    TODO: Improve accuracy by using a more sophisticated method, such as a machine learning model.
    """
    # Split content into words and convert to lowercase
    words = file_content.lower().split()
    
    # Check each word against industry keywords
    for word in words:
        # Check civil keywords
        if word in civil_classifier.CIVIL_KEYWORDS:
            return "civil"
            
        # Check finance keywords 
        if word in finance_classifier.FINANCE_KEYWORDS:
            return "finance"
            
    # No industry keywords found
    return "unknown"

def classify_file(file: FileStorage) -> str:
    """
    Classifies a file by first determining its industry and then using the appropriate industry-specific classifier.

    Args:
        file (FileStorage): The file to classify, provided as a FileStorage object

    Returns:
        str: The predicted class of the file. Returns 'unknown' if:
            - The industry cannot be determined
            - The industry is not supported
            - The file cannot be classified by the industry-specific classifier

    The function follows these steps:
    1. Extracts text content from the file
    2. Determines the industry based on keywords
    3. Uses the corresponding industry classifier to determine the specific document type
    """
    # Classifies the file based on the content
    file_content = get_file_content(file)
    # Classifies the industry of the file
    industry = classify_industry(file_content)
    if industry not in SUPPORTED_INDUSTRIES:
        return "unknown"
    # Uses the corresponding industry-specific classifier
    if industry == "civil":
        classifier = civil_classifier.CivilClassifier()
    elif industry == "finance":
        classifier = finance_classifier.FinanceClassifier()
    else: 
        return "unknown"
    return classifier.classify(file_content)

